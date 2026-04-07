#!/usr/bin/env python3
"""Generate an AI training lesson plan PPTX with skills introduction."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent
OUTPUT_DIR = PROJECT_ROOT / "output"
OUTPUT_FILE = OUTPUT_DIR / "AI教程教案_Skills介绍_20260405.pptx"

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"

COLORS = {
    "navy": RGBColor(20, 37, 63),
    "ink": RGBColor(31, 41, 55),
    "muted": RGBColor(99, 115, 129),
    "sand": RGBColor(247, 242, 234),
    "cream": RGBColor(255, 251, 245),
    "gold": RGBColor(212, 165, 116),
    "coral": RGBColor(228, 107, 92),
    "teal": RGBColor(52, 130, 126),
    "sage": RGBColor(160, 182, 150),
    "sky": RGBColor(210, 226, 244),
    "white": RGBColor(255, 255, 255),
    "line": RGBColor(220, 212, 201),
    "soft_navy": RGBColor(231, 237, 247),
}


def set_shape_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_line(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = FONT_CN,
    margin: float = 0.08,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_multiline_text(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    *,
    font_size: int = 16,
    color: RGBColor | None = None,
    bold_first: bool = False,
    line_spacing: float = 1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(3)
        paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT_CN
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color or COLORS["ink"]
    return box


def add_tag(slide, left, top, text: str, bg: RGBColor, fg: RGBColor | None = None):
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        Inches(1.15),
        Inches(0.38),
    )
    set_shape_fill(tag, bg)
    tag.line.fill.background()
    add_textbox(
        slide,
        left,
        top + Inches(0.02),
        Inches(1.15),
        Inches(0.3),
        text,
        font_size=10,
        color=fg or COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_title(slide, number: int, title: str, subtitle: str | None = None) -> None:
    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.45),
        Inches(9.7),
        Inches(0.55),
        title,
        font_size=27,
        color=COLORS["navy"],
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.68),
            Inches(0.92),
            Inches(10.5),
            Inches(0.38),
            subtitle,
            font_size=11,
            color=COLORS["muted"],
        )
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(11.95),
        Inches(0.42),
        Inches(0.7),
        Inches(0.42),
    )
    set_shape_fill(badge, COLORS["navy"])
    badge.line.fill.background()
    add_textbox(
        slide,
        Inches(11.95),
        Inches(0.44),
        Inches(0.7),
        Inches(0.3),
        f"{number:02d}",
        font_size=12,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_note_box(slide, text: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(6.62),
        Inches(12.0),
        Inches(0.48),
    )
    set_shape_fill(shape, COLORS["soft_navy"])
    shape.line.fill.background()
    add_textbox(
        slide,
        Inches(0.78),
        Inches(6.68),
        Inches(11.8),
        Inches(0.26),
        f"讲师提示：{text}",
        font_size=10,
        color=COLORS["navy"],
        bold=True,
    )


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    body_lines: list[str],
    *,
    fill: RGBColor,
    title_color: RGBColor,
    body_color: RGBColor | None = None,
    accent: RGBColor | None = None,
):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    set_shape_fill(card, fill)
    set_shape_line(card, COLORS["line"], 1)
    if accent:
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left,
            top,
            width,
            Inches(0.08),
        )
        set_shape_fill(bar, accent)
        bar.line.fill.background()
    add_textbox(
        slide,
        left + Inches(0.14),
        top + Inches(0.16),
        width - Inches(0.28),
        Inches(0.34),
        title,
        font_size=15,
        color=title_color,
        bold=True,
    )
    add_multiline_text(
        slide,
        left + Inches(0.14),
        top + Inches(0.52),
        width - Inches(0.28),
        height - Inches(0.62),
        body_lines,
        font_size=11,
        color=body_color or COLORS["ink"],
    )


def add_flow_arrow(slide, x1, y1, x2, y2, color: RGBColor) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.5)
    line.line.end_arrowhead = True


def add_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(bg, COLORS["navy"])
    bg.line.fill.background()

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.9), Inches(-0.5), Inches(5.5), Inches(5.5))
    set_shape_fill(glow, COLORS["teal"])
    glow.fill.transparency = 0.55
    glow.line.fill.background()

    glow2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-1.2), Inches(4.6), Inches(4.4), Inches(4.4))
    set_shape_fill(glow2, COLORS["coral"])
    glow2.fill.transparency = 0.68
    glow2.line.fill.background()

    add_tag(slide, Inches(0.72), Inches(0.72), "PPTX", COLORS["gold"], COLORS["navy"])
    add_textbox(
        slide,
        Inches(0.72),
        Inches(1.3),
        Inches(8.2),
        Inches(1.7),
        "AI 教程教案",
        font_size=31,
        color=COLORS["white"],
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.74),
        Inches(2.4),
        Inches(8.9),
        Inches(0.9),
        "用 `pptx` 技能做一份能讲、能改、能复用的培训演示稿",
        font_size=16,
        color=COLORS["sand"],
    )

    for idx, text in enumerate(
        [
            "适用对象：零基础或轻度接触过 AI 的职场人",
            "课时建议：30 到 40 分钟",
            "本课重点：AI 入门认知 + Skills 介绍 + 提问模板",
        ]
    ):
        left = Inches(0.75 + idx * 4.1)
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(5.45),
            Inches(3.7),
            Inches(0.58),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLORS["white"]
        pill.fill.transparency = 0.84
        pill.line.fill.background()
        add_textbox(
            slide,
            left + Inches(0.08),
            Inches(5.57),
            Inches(3.5),
            Inches(0.2),
            text,
            font_size=10,
            color=COLORS["white"],
            bold=True,
        )

    add_textbox(
        slide,
        Inches(0.76),
        Inches(6.45),
        Inches(4.8),
        Inches(0.28),
        "输出示例：这份 PPT 本身即为 skills 生成教案范例",
        font_size=10,
        color=COLORS["sand"],
    )


def add_slide_02(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 2, "这节课的目标", "让听众在一节课内把 AI、Skills 和落地用法串起来")
    cards = [
        ("认知归位", ["分清公司、模型、产品、Skill、Agent", "不再把名字混在一起"], COLORS["cream"], COLORS["navy"], COLORS["gold"]),
        ("场景入门", ["先抓住问答、写作、办公整理 3 类高频任务", "避免一上来追求复杂自动化"], COLORS["sand"], COLORS["navy"], COLORS["coral"]),
        ("认识 Skills", ["知道为什么 Skill 能让 AI 更稳定地产出文件", "理解办公四件套和组合用法"], COLORS["soft_navy"], COLORS["navy"], COLORS["teal"]),
        ("马上能练", ["带走一个 4 步提问模板", "回去就能拿真实任务测试"], COLORS["white"], COLORS["navy"], COLORS["sage"]),
    ]
    positions = [
        (Inches(0.7), Inches(1.55)),
        (Inches(6.75), Inches(1.55)),
        (Inches(0.7), Inches(3.75)),
        (Inches(6.75), Inches(3.75)),
    ]
    for (title, lines, fill, color, accent), (left, top) in zip(cards, positions):
        add_card(slide, left, top, Inches(5.8), Inches(1.8), title, lines, fill=fill, title_color=color, accent=accent)
    add_note_box(slide, "这一页先帮听众降压，强调不是开发课，而是帮大家建立一套更清晰的使用框架。")


def add_slide_03(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 3, "先建立一个最重要的认知", "不是所有 AI 名字都在同一层；名字归位，后面才不容易混乱")
    stages = [
        ("公司", "OpenAI / Google / Anthropic", COLORS["navy"], COLORS["white"]),
        ("模型", "GPT / Gemini / Claude", COLORS["gold"], COLORS["navy"]),
        ("产品", "ChatGPT / Gemini App / Claude", COLORS["teal"], COLORS["white"]),
        ("Skill", "pptx / pdf / xlsx / frontend-design", COLORS["coral"], COLORS["white"]),
        ("Agent", "OpenClaw / 自动化工作流", COLORS["sage"], COLORS["ink"]),
    ]
    lefts = [0.75, 3.2, 5.65, 8.1, 10.55]
    for idx, ((title, body, fill, fg), x) in enumerate(zip(stages, lefts)):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(2.2),
            Inches(2.0),
            Inches(2.05),
        )
        set_shape_fill(card, fill)
        card.line.fill.background()
        add_textbox(
            slide,
            Inches(x + 0.15),
            Inches(2.42),
            Inches(1.7),
            Inches(0.32),
            title,
            font_size=16,
            color=fg,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_multiline_text(
            slide,
            Inches(x + 0.12),
            Inches(2.9),
            Inches(1.76),
            Inches(0.9),
            [body],
            font_size=10,
            color=fg,
        )
        if idx < len(stages) - 1:
            add_flow_arrow(
                slide,
                Inches(x + 2.0),
                Inches(3.22),
                Inches(lefts[idx + 1]),
                Inches(3.22),
                COLORS["muted"],
            )
    add_textbox(
        slide,
        Inches(1.25),
        Inches(4.75),
        Inches(10.9),
        Inches(0.65),
        "一句话总结：公司造模型，模型装进产品，Skill 让 AI 学会做具体文件任务，Agent 再把这些能力接进真实流程。",
        font_size=15,
        color=COLORS["navy"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_note_box(slide, "这里建议停顿一下，用例子解释：ChatGPT 是产品，GPT 是模型，pptx 是让 AI 处理演示文稿的专业技能。")


def add_slide_04(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 4, "新手最常见的 3 个 AI 使用场景", "先从高频低门槛场景起步，比一开始追求全自动更容易建立信心")
    blocks = [
        ("01", "问答与学习", ["解释概念", "做入门计划", "总结一段材料"], COLORS["soft_navy"], COLORS["navy"]),
        ("02", "内容与写作", ["写邮件文案", "改标题语气", "生成大纲初稿"], COLORS["cream"], COLORS["coral"]),
        ("03", "办公与整理", ["会议纪要", "文档重点提炼", "表格与汇报准备"], COLORS["sand"], COLORS["teal"]),
    ]
    for idx, (num, title, lines, fill, accent) in enumerate(blocks):
        left = Inches(0.85 + idx * 4.15)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.75), Inches(3.55), Inches(3.85))
        set_shape_fill(card, fill)
        set_shape_line(card, COLORS["line"], 1)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), Inches(2.0), Inches(0.7), Inches(0.7))
        set_shape_fill(circle, accent)
        circle.line.fill.background()
        add_textbox(slide, left + Inches(0.18), Inches(2.15), Inches(0.7), Inches(0.22), num, font_size=14, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_EN)
        add_textbox(slide, left + Inches(1.02), Inches(2.08), Inches(2.0), Inches(0.3), title, font_size=16, color=COLORS["ink"], bold=True)
        add_multiline_text(slide, left + Inches(0.22), Inches(2.9), Inches(3.08), Inches(1.5), [f"• {line}" for line in lines], font_size=12, color=COLORS["ink"])
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(5.95), Inches(5.0), Inches(0.46))
    set_shape_fill(ribbon, COLORS["navy"])
    ribbon.line.fill.background()
    add_textbox(slide, Inches(4.12), Inches(6.03), Inches(4.95), Inches(0.22), "这 3 类通常能覆盖大多数普通用户 80% 的首次需求", font_size=11, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    add_note_box(slide, "建议用真实例子带一下，比如写邮件、整理长文、做简报提纲，让听众立刻知道 AI 能替自己省哪类时间。")


def add_slide_05(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 5, "为什么要讲 Skills", "Skill 可以理解成给 AI 的专业操作说明书，让结果更稳、更像成品")

    add_card(
        slide,
        Inches(0.8),
        Inches(1.7),
        Inches(5.65),
        Inches(3.9),
        "只有聊天窗口时",
        ["提问后拿到一段答案", "还要自己搬到 Word / Excel / PPT", "格式容易散，重复步骤很多", "一换任务就要重新摸索"],
        fill=COLORS["cream"],
        title_color=COLORS["navy"],
        accent=COLORS["coral"],
    )
    add_card(
        slide,
        Inches(6.9),
        Inches(1.7),
        Inches(5.65),
        Inches(3.9),
        "接入 Skill 之后",
        ["AI 知道该如何处理具体文件格式", "可以直接创建、读取、编辑交付物", "路径更清晰，复用更容易", "从“给建议”更接近“帮你产出”"],
        fill=COLORS["soft_navy"],
        title_color=COLORS["navy"],
        accent=COLORS["teal"],
    )
    add_flow_arrow(slide, Inches(5.7), Inches(3.63), Inches(6.72), Inches(3.63), COLORS["gold"])
    add_textbox(slide, Inches(5.38), Inches(3.28), Inches(1.65), Inches(0.28), "加入 Skill", font_size=11, color=COLORS["gold"], bold=True, align=PP_ALIGN.CENTER)
    add_note_box(slide, "这一页重点不是讲技术细节，而是让听众明白：Skill 的价值在于把零散回答变成更稳定的文件产出流程。")


def add_slide_06(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 6, "办公四件套为什么值得优先了解", "如果你的工作经常写总结、做表格、整理资料、做汇报，这组通常是第一梯队")
    skills = [
        ("docx", "写报告、方案、通知、合同草稿", COLORS["soft_navy"], COLORS["navy"]),
        ("xlsx", "做台账、整理数据、批量填充结构化表格", COLORS["cream"], COLORS["teal"]),
        ("pdf", "读资料、抽重点、提取文本和表格", COLORS["sand"], COLORS["coral"]),
        ("pptx", "做课件、汇报、提案和演示稿", COLORS["white"], COLORS["gold"]),
    ]
    positions = [(0.9, 1.75), (6.85, 1.75), (0.9, 3.9), (6.85, 3.9)]
    for (name, desc, fill, accent), (x, y) in zip(skills, positions):
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.55), Inches(1.7))
        set_shape_fill(card, fill)
        set_shape_line(card, COLORS["line"], 1)
        icon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.28), Inches(0.58), Inches(0.58))
        set_shape_fill(icon, accent)
        icon.line.fill.background()
        add_textbox(slide, Inches(x + 0.9), Inches(y + 0.3), Inches(2.0), Inches(0.25), name, font_size=18, color=COLORS["navy"], bold=True, font_name=FONT_EN)
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.92), Inches(4.95), Inches(0.4), desc, font_size=11, color=COLORS["ink"])
    add_note_box(slide, "不要把它们理解成“只是支持某个格式”，更准确地说，是让 AI 围绕这些文件类型工作得更稳定。")


def add_slide_07(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 7, "重点示例：`pptx` skill 如何帮你做教案", "今天这份教案就属于典型的 `pptx` 使用场景")
    steps = [
        ("1", "定目标", "先确定听众是谁、课时多久、目标是什么"),
        ("2", "排页纲", "把每页的主题、要点、讲师提示先写清楚"),
        ("3", "出文件", "生成可继续修改的 `.pptx` 成品"),
        ("4", "复改稿", "再根据课堂反馈增删页面和措辞"),
    ]
    for idx, (num, title, desc) in enumerate(steps):
        x = 0.95 + idx * 3.05
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.1), Inches(0.78), Inches(0.78))
        set_shape_fill(circle, COLORS["navy"] if idx % 2 == 0 else COLORS["teal"])
        circle.line.fill.background()
        add_textbox(slide, Inches(x), Inches(2.27), Inches(0.78), Inches(0.2), num, font_size=15, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_EN)
        add_textbox(slide, Inches(x - 0.05), Inches(3.0), Inches(1.2), Inches(0.28), title, font_size=14, color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x - 0.42), Inches(3.43), Inches(1.95), Inches(0.7), desc, font_size=10, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_flow_arrow(slide, Inches(x + 0.86), Inches(2.49), Inches(x + 2.22), Inches(2.49), COLORS["gold"])
    add_card(
        slide,
        Inches(9.25),
        Inches(4.05),
        Inches(3.15),
        Inches(1.7),
        "特别适合的任务",
        ["培训课件", "项目汇报", "销售提案", "活动路演稿"],
        fill=COLORS["cream"],
        title_color=COLORS["navy"],
        accent=COLORS["gold"],
    )
    add_note_box(slide, "这里可以顺手提醒：不是先做漂亮，而是先把页纲和每页目标讲清楚，PPT 的质量会稳定很多。")


def add_slide_08(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 8, "更值得优先了解的 6 类 Skills", "不是要求一次装满，而是先找最贴近自己产出的那几类")
    items = [
        ("Frontend Design", "先把页面或演示稿的审美方向想清楚", COLORS["soft_navy"], COLORS["navy"]),
        ("办公四件套", "docx / xlsx / pdf / pptx，围绕文件稳定产出", COLORS["cream"], COLORS["coral"]),
        ("Web Access", "带着浏览器上下文去搜索、进入页面、执行操作", COLORS["sand"], COLORS["teal"]),
        ("PUA", "AI 卡住或摆烂时，用来打断惯性思路", COLORS["white"], COLORS["navy"]),
        ("Claude-mem", "给长期项目补上一层跨会话记忆", COLORS["soft_navy"], COLORS["teal"]),
        ("Skill-Creator", "把高频工作流沉淀成自己的 Skill", COLORS["cream"], COLORS["gold"]),
    ]
    for idx, (title, desc, fill, accent) in enumerate(items):
        col = idx % 3
        row = idx // 3
        x = 0.75 + col * 4.15
        y = 1.75 + row * 2.05
        add_card(slide, Inches(x), Inches(y), Inches(3.7), Inches(1.65), title, [desc], fill=fill, title_color=COLORS["navy"], body_color=COLORS["ink"], accent=accent)
    add_note_box(slide, "如果时间有限，就重点讲 Frontend Design、办公四件套和 Skill-Creator，这三类最容易和听众已有工作连接起来。")


def add_slide_09(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 9, "组合起来用，效果往往比单装一个更好", "真正有价值的不是某一个 Skill，而是把几类能力串成一条工作流")
    combos = [
        ("组合 1", "Frontend Design", "pptx / Office 文档工具", "先定审美方向，再生成和细修演示稿"),
        ("组合 2", "pdf", "docx / xlsx / pptx", "先读资料提炼重点，再转换成交付文件"),
        ("组合 3", "Web Access", "Claude-mem", "一个负责搜和操作，一个负责长期记忆"),
    ]
    for idx, (tag, left_skill, right_skill, desc) in enumerate(combos):
        y = 1.85 + idx * 1.48
        add_tag(slide, Inches(0.88), Inches(y), tag, COLORS["navy"])
        left_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(y - 0.02), Inches(2.6), Inches(0.56))
        set_shape_fill(left_box, COLORS["soft_navy"])
        set_shape_line(left_box, COLORS["line"], 1)
        add_textbox(slide, Inches(2.25), Inches(y + 0.09), Inches(2.45), Inches(0.22), left_skill, font_size=13, color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)
        add_flow_arrow(slide, Inches(4.88), Inches(y + 0.25), Inches(6.15), Inches(y + 0.25), COLORS["gold"])
        right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.25), Inches(y - 0.02), Inches(2.9), Inches(0.56))
        set_shape_fill(right_box, COLORS["cream"])
        set_shape_line(right_box, COLORS["line"], 1)
        add_textbox(slide, Inches(6.32), Inches(y + 0.09), Inches(2.75), Inches(0.22), right_skill, font_size=13, color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(9.45), Inches(y + 0.01), Inches(2.8), Inches(0.46), desc, font_size=10, color=COLORS["muted"])
    add_note_box(slide, "这一页适合帮大家建立“组合思维”：先研究、再生成、再沉淀，比把所有任务都扔给同一个窗口更稳。")


def add_slide_10(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 10, "一套最适合新手的 4 步提问法", "很多人觉得 AI 不好用，不是模型不行，而是提问还停留在碰运气")
    steps = [
        ("角色", "你希望 AI 扮演谁"),
        ("任务", "你要它完成什么"),
        ("背景", "对象、场景、限制条件"),
        ("输出", "你要什么格式和风格"),
    ]
    for idx, (title, desc) in enumerate(steps):
        x = 0.85 + idx * 3.05
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.55), Inches(1.45))
        set_shape_fill(card, COLORS["cream"] if idx % 2 == 0 else COLORS["soft_navy"])
        set_shape_line(card, COLORS["line"], 1)
        add_textbox(slide, Inches(x + 0.18), Inches(2.03), Inches(2.15), Inches(0.25), title, font_size=16, color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.15), Inches(2.5), Inches(2.2), Inches(0.35), desc, font_size=11, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    sample = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(3.7), Inches(11.5), Inches(2.15))
    set_shape_fill(sample, COLORS["navy"])
    sample.line.fill.background()
    add_textbox(slide, Inches(1.1), Inches(3.95), Inches(11.0), Inches(0.28), "示例提示词", font_size=14, color=COLORS["gold"], bold=True)
    add_multiline_text(
        slide,
        Inches(1.08),
        Inches(4.35),
        Inches(10.9),
        Inches(1.2),
        [
            "你现在是一名资深培训讲师。",
            "请帮我做一份 AI 入门课 PPT，听众是零基础职场人，时长 30 分钟。",
            "请按“每页标题 + 3 个要点 + 讲师提示”的格式输出，并补一页 Skills 介绍。",
        ],
        font_size=13,
        color=COLORS["white"],
    )
    add_note_box(slide, "这页适合现场直接演示：把模糊提问和结构化提问对比一下，听众会立刻感受到差别。")


def add_slide_11(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 11, "这份教案可以怎么讲", "如果你拿去做一次 30 分钟分享，可以按下面节奏展开")
    headers = ["时间", "讲什么", "讲师动作", "听众收获"]
    col_x = [0.85, 2.05, 5.15, 8.85]
    col_w = [1.0, 2.8, 3.35, 3.45]
    for x, w, head in zip(col_x, col_w, headers):
        head_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.62), Inches(w), Inches(0.48))
        set_shape_fill(head_box, COLORS["navy"])
        head_box.line.fill.background()
        add_textbox(slide, Inches(x), Inches(1.72), Inches(w), Inches(0.2), head, font_size=12, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
    rows = [
        ("5 分钟", "先讲 AI 名字怎么区分", "用“公司-模型-产品-Skill-Agent”框架举例", "概念不再混乱"),
        ("8 分钟", "讲 3 个高频使用场景", "展示写邮件、总结文档、做提纲", "知道 AI 先帮自己省哪种时间"),
        ("10 分钟", "讲 Skills 与办公四件套", "重点放在 pptx、pdf、docx、xlsx 的价值", "知道 Skill 为什么能提高稳定性"),
        ("7 分钟", "做一次结构化提问演示", "现场用 4 步提问法生成一个小任务", "带走能马上复用的模板"),
    ]
    for idx, row in enumerate(rows):
        y = 2.25 + idx * 1.0
        for jdx, (x, w, text) in enumerate(zip(col_x, col_w, row)):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.8))
            set_shape_fill(cell, COLORS["cream"] if idx % 2 == 0 else COLORS["soft_navy"])
            set_shape_line(cell, COLORS["line"], 0.8)
            add_textbox(
                slide,
                Inches(x + 0.06),
                Inches(y + 0.12),
                Inches(w - 0.12),
                Inches(0.5),
                text,
                font_size=10 if jdx != 0 else 11,
                color=COLORS["ink"],
                bold=jdx == 0,
                align=PP_ALIGN.CENTER if jdx == 0 else PP_ALIGN.LEFT,
            )
    add_note_box(slide, "如果时间不够，就把技能介绍压缩成办公四件套 + 一个组合案例，最后把提问模板完整讲透。")


def add_slide_12(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, 12, "结尾页：让听众带走什么", "一节入门课最重要的不是信息量，而是让大家敢回去自己试一次")
    takeaways = [
        ("先选 1 个主力 AI", "别一开始同时学 5 个工具，先有主力，再补备选"),
        ("再补 1 个 Skill", "如果经常做文档和汇报，优先从办公四件套开始"),
        ("拿 1 个真实任务练", "用写邮件、做提纲、整理资料这种手头任务试一次"),
    ]
    for idx, (title, desc) in enumerate(takeaways):
        left = Inches(0.95 + idx * 4.05)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(3.45), Inches(2.6))
        set_shape_fill(card, COLORS["navy"] if idx == 1 else COLORS["cream"])
        set_shape_line(card, COLORS["line"], 1)
        fg = COLORS["white"] if idx == 1 else COLORS["navy"]
        add_textbox(slide, left + Inches(0.18), Inches(2.32), Inches(3.05), Inches(0.34), title, font_size=17, color=fg, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.18), Inches(3.02), Inches(3.05), Inches(0.82), desc, font_size=11, color=fg, align=PP_ALIGN.CENTER)
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.55), Inches(5.25), Inches(6.15), Inches(0.62))
    set_shape_fill(footer, COLORS["gold"])
    footer.line.fill.background()
    add_textbox(slide, Inches(3.6), Inches(5.43), Inches(6.0), Inches(0.22), "建议行动：今晚就用一个真实工作任务，试一次“AI + Skill”的组合", font_size=12, color=COLORS["navy"], bold=True, align=PP_ALIGN.CENTER)
    add_note_box(slide, "结尾不要再堆概念，用一句明确行动建议收束，让大家课后立刻去试。")


def build_presentation() -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs)
    add_slide_02(prs)
    add_slide_03(prs)
    add_slide_04(prs)
    add_slide_05(prs)
    add_slide_06(prs)
    add_slide_07(prs)
    add_slide_08(prs)
    add_slide_09(prs)
    add_slide_10(prs)
    add_slide_11(prs)
    add_slide_12(prs)

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


if __name__ == "__main__":
    output = build_presentation()
    print(output)
